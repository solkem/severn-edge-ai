"""Generate a high-quality, editable 5-slide classroom deck."""

from __future__ import annotations

import math
import os
import random

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = os.path.join(os.path.dirname(__file__), "Severn_Edge_AI_Classroom_Slides.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "navy": (7, 24, 61),
    "blue": (25, 76, 173),
    "cyan": (77, 219, 255),
    "teal": (36, 205, 167),
    "green": (76, 212, 123),
    "yellow": (255, 212, 74),
    "orange": (255, 152, 66),
    "pink": (255, 90, 166),
    "purple": (118, 92, 255),
    "white": (248, 251, 255),
    "muted": (184, 203, 234),
    "panel": (14, 41, 91),
    "panel2": (13, 31, 70),
}


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def clamp(value: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, value))


def fill_solid(shape, color: tuple[int, int, int], transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = clamp(transparency, 0.0, 1.0)


def set_line(shape, color: tuple[int, int, int], width_pt: float = 1.0, transparency: float = 0.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width_pt)
    shape.line.transparency = clamp(transparency, 0.0, 1.0)


def hide_line(shape) -> None:
    shape.line.fill.background()


def add_bg(slide, color: tuple[int, int, int]) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_solid(rect, color)
    hide_line(rect)


def add_glow(slide, left, top, width, height, color: tuple[int, int, int], alpha: float) -> None:
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    fill_solid(glow, color, alpha)
    hide_line(glow)


def add_textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    size_pt: float = 24,
    bold: bool = False,
    color: tuple[int, int, int] = COLORS["white"],
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
    vertical_anchor=MSO_ANCHOR.MIDDLE,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if not p.runs:
            continue
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_card(
    slide,
    left,
    top,
    width,
    height,
    fill: tuple[int, int, int] = COLORS["panel"],
    border: tuple[int, int, int] = COLORS["blue"],
    radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    card = slide.shapes.add_shape(radius_shape, left, top, width, height)
    fill_solid(card, fill)
    set_line(card, border, width_pt=1.25)
    return card


def add_note(slide, text: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = text


def draw_polyline(slide, points, color: tuple[int, int, int], width_pt: float = 1.75) -> None:
    for i in range(len(points) - 1):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
        line.line.color.rgb = rgb(color)
        line.line.width = Pt(width_pt)


def interpolate(c1: tuple[int, int, int], c2: tuple[int, int, int], t: float) -> tuple[int, int, int]:
    t = clamp(t, 0.0, 1.0)
    return (
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t),
    )


def heat_color(v: float) -> tuple[int, int, int]:
    stops = [
        (-1.0, (29, 78, 216)),
        (-0.3, (71, 211, 255)),
        (0.0, (228, 243, 117)),
        (0.45, (255, 169, 61)),
        (1.0, (222, 66, 72)),
    ]
    for i in range(len(stops) - 1):
        x1, c1 = stops[i]
        x2, c2 = stops[i + 1]
        if x1 <= v <= x2:
            return interpolate(c1, c2, (v - x1) / (x2 - x1))
    return stops[0][1] if v < stops[0][0] else stops[-1][1]


def draw_heatmap(slide, left, top, width, height, rows: int, cols: int, value_fn) -> None:
    cell_w = max(1, int(width / cols))
    cell_h = max(1, int(height / rows))
    for r in range(rows):
        for c in range(cols):
            v = value_fn(r, c)
            color = heat_color(v)
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(left + c * cell_w),
                int(top + r * cell_h),
                cell_w + 1,
                cell_h + 1,
            )
            fill_solid(cell, color)
            hide_line(cell)


def build_slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["navy"])
    add_glow(slide, Inches(3.6), Inches(1.35), Inches(6.1), Inches(3.25), COLORS["blue"], 0.65)
    add_glow(slide, Inches(4.3), Inches(1.8), Inches(4.8), Inches(2.4), COLORS["cyan"], 0.76)

    add_textbox(
        slide,
        "SEVERN EDGE AI",
        Inches(0.6),
        Inches(0.35),
        Inches(12.1),
        Inches(0.85),
        size_pt=56,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        "Today, YOU teach a robot",
        Inches(0.8),
        Inches(1.15),
        Inches(11.7),
        Inches(0.65),
        size_pt=30,
        color=COLORS["cyan"],
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    chip = add_card(slide, Inches(4.32), Inches(2.02), Inches(4.7), Inches(2.15), fill=COLORS["panel2"], border=COLORS["cyan"])
    chip.adjustments[0] = 0.08
    inner = add_card(slide, Inches(4.75), Inches(2.43), Inches(3.84), Inches(1.34), fill=(21, 69, 122), border=COLORS["teal"])
    inner.adjustments[0] = 0.05
    add_textbox(
        slide,
        "Arduino Nano 33 BLE",
        Inches(4.95),
        Inches(2.77),
        Inches(3.45),
        Inches(0.44),
        size_pt=20,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        "Your board learns gestures from your data",
        Inches(4.85),
        Inches(3.2),
        Inches(3.65),
        Inches(0.35),
        size_pt=16,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )

    draw_polyline(
        slide,
        [
            (Inches(4.25), Inches(2.55)),
            (Inches(3.55), Inches(2.22)),
            (Inches(2.75), Inches(2.22)),
        ],
        COLORS["cyan"],
    )
    draw_polyline(
        slide,
        [
            (Inches(9.06), Inches(2.55)),
            (Inches(9.77), Inches(2.22)),
            (Inches(10.53), Inches(2.22)),
        ],
        COLORS["cyan"],
    )
    draw_polyline(
        slide,
        [
            (Inches(4.25), Inches(3.65)),
            (Inches(3.52), Inches(3.98)),
            (Inches(2.75), Inches(3.98)),
        ],
        COLORS["cyan"],
    )
    draw_polyline(
        slide,
        [
            (Inches(9.06), Inches(3.65)),
            (Inches(9.76), Inches(3.98)),
            (Inches(10.53), Inches(3.98)),
        ],
        COLORS["cyan"],
    )

    steps = ["Connect", "Collect", "Train", "Test"]
    x_start = Inches(1.02)
    card_w = Inches(2.9)
    gap = Inches(0.28)
    for idx, step in enumerate(steps):
        x = int(x_start + idx * (card_w + gap))
        add_card(slide, x, Inches(5.7), card_w, Inches(1.0), fill=(14, 50, 108), border=COLORS["blue"])
        add_textbox(slide, step, x, Inches(5.7), card_w, Inches(1.0), size_pt=28, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, int(x + card_w + Inches(0.07)), Inches(6.0), Inches(0.16), Inches(0.42))
            fill_solid(arrow, COLORS["cyan"])
            hide_line(arrow)

    add_note(
        slide,
        "Hook: 'Is AI magic?'\n"
        "This board will learn from students today.\n"
        "Set expectation for workflow: Connect -> Collect -> Train -> Test.",
    )


def build_slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["navy"])
    add_glow(slide, Inches(2.9), Inches(0.9), Inches(7.6), Inches(4.2), COLORS["blue"], 0.72)

    add_textbox(
        slide,
        "How the Sensor Captures Your Move",
        Inches(0.55),
        Inches(0.28),
        Inches(12.25),
        Inches(0.8),
        size_pt=46,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Display",
    )

    board = add_card(slide, Inches(3.92), Inches(1.5), Inches(5.45), Inches(2.65), fill=(235, 241, 248), border=(165, 183, 205))
    board.adjustments[0] = 0.04
    add_card(slide, Inches(4.2), Inches(1.77), Inches(4.88), Inches(2.1), fill=(246, 250, 255), border=(201, 217, 235))
    add_textbox(slide, "IMU sensor", Inches(5.85), Inches(2.33), Inches(1.7), Inches(0.4), size_pt=16, bold=True, color=(40, 80, 120), align=PP_ALIGN.CENTER)

    # Accelerometer axes (red)
    ax_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.68), Inches(2.75), Inches(2.1), Inches(0.33))
    fill_solid(ax_arrow, (246, 84, 84))
    hide_line(ax_arrow)
    add_textbox(slide, "ax (left-right)", Inches(8.77), Inches(2.66), Inches(2.35), Inches(0.5), size_pt=22, color=COLORS["white"])

    az_arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(6.22), Inches(1.25), Inches(0.4), Inches(1.6))
    fill_solid(az_arrow, (246, 84, 84))
    hide_line(az_arrow)
    add_textbox(slide, "az (up-down)", Inches(5.35), Inches(1.05), Inches(2.1), Inches(0.36), size_pt=22, align=PP_ALIGN.CENTER)

    ay_arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.78), Inches(3.05), Inches(1.8), Inches(0.32))
    ay_arrow.rotation = 145
    fill_solid(ay_arrow, (246, 84, 84))
    hide_line(ay_arrow)
    add_textbox(slide, "ay (forward-back)", Inches(1.82), Inches(3.83), Inches(2.55), Inches(0.5), size_pt=22)

    # Gyroscope labels (blue)
    add_textbox(slide, "gx (rotation)", Inches(8.78), Inches(2.22), Inches(2.1), Inches(0.4), size_pt=20, color=COLORS["cyan"])
    add_textbox(slide, "gy (rotation)", Inches(1.82), Inches(3.22), Inches(2.0), Inches(0.4), size_pt=20, color=COLORS["cyan"])
    add_textbox(slide, "gz (rotation)", Inches(6.9), Inches(1.52), Inches(1.8), Inches(0.4), size_pt=20, color=COLORS["cyan"])
    ring1 = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(8.17), Inches(2.13), Inches(0.55), Inches(0.55))
    fill_solid(ring1, COLORS["cyan"], 0.75)
    set_line(ring1, COLORS["cyan"], 1.0)
    ring2 = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(2.85), Inches(3.17), Inches(0.55), Inches(0.55))
    fill_solid(ring2, COLORS["cyan"], 0.75)
    set_line(ring2, COLORS["cyan"], 1.0)
    ring3 = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(6.3), Inches(1.53), Inches(0.52), Inches(0.52))
    fill_solid(ring3, COLORS["cyan"], 0.75)
    set_line(ring3, COLORS["cyan"], 1.0)

    # Six readable signal cards in a 3x2 grid (clear from back of classroom)
    axes = [
        ("ax", "left-right", COLORS["orange"]),
        ("ay", "forward-back", COLORS["orange"]),
        ("az", "up-down", COLORS["orange"]),
        ("gx", "rotation", COLORS["cyan"]),
        ("gy", "rotation", COLORS["cyan"]),
        ("gz", "rotation", COLORS["cyan"]),
    ]
    start_x = Inches(0.42)
    start_y = Inches(4.74)
    card_w = Inches(4.1)
    card_h = Inches(1.32)
    gap_x = Inches(0.2)
    gap_y = Inches(0.14)
    for idx, (axis, label, line_color) in enumerate(axes):
        row = idx // 3
        col = idx % 3
        left = int(start_x + col * (card_w + gap_x))
        top = int(start_y + row * (card_h + gap_y))
        add_card(slide, left, top, card_w, card_h, fill=COLORS["panel2"], border=(42, 108, 195))
        add_textbox(
            slide,
            f"{axis} ({label})",
            left + Inches(0.1),
            top + Inches(0.06),
            card_w - Inches(0.2),
            Inches(0.3),
            size_pt=20,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        chart_left = int(left + Inches(0.22))
        chart_top = int(top + Inches(0.46))
        chart_w = int(card_w - Inches(0.44))
        chart_h = int(card_h - Inches(0.56))
        baseline = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            chart_left,
            chart_top + chart_h // 2,
            chart_left + chart_w,
            chart_top + chart_h // 2,
        )
        baseline.line.color.rgb = rgb((126, 168, 220))
        baseline.line.width = Pt(1.0)

        points = []
        for s in range(14):
            t = s / 13
            amp = 0.34 if idx < 3 else 0.24
            val = math.sin(t * math.pi * 2 + idx * 0.52) * amp + (0.05 if idx == 0 else 0)
            x = chart_left + int(t * chart_w)
            y = chart_top + int(chart_h * (0.5 - val))
            points.append((x, y))
        draw_polyline(slide, points, line_color, width_pt=2.0)

    add_note(
        slide,
        "Students see 6 sensor channels in real time.\n"
        "Red axes are acceleration (ax, ay, az); blue axes are rotation (gx, gy, gz).\n"
        "Main teaching line: AI does not see motion directly, only these changing numbers.\n"
        "Use the 3x2 card grid to ask: which channels change most for each movement?",
    )


def build_slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["navy"])
    add_glow(slide, Inches(2.2), Inches(1.0), Inches(8.8), Inches(4.3), COLORS["purple"], 0.78)

    add_textbox(
        slide,
        "What the AI Actually Sees",
        Inches(0.55),
        Inches(0.23),
        Inches(12.25),
        Inches(0.8),
        size_pt=44,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        "100 time steps x 6 axes = 600 numbers",
        Inches(3.0),
        Inches(0.92),
        Inches(7.4),
        Inches(0.4),
        size_pt=24,
        color=COLORS["cyan"],
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    panel_left = Inches(1.0)
    panel_top = Inches(1.35)
    panel_w = Inches(11.3)
    panel_h = Inches(3.35)
    add_card(slide, panel_left, panel_top, panel_w, panel_h, fill=COLORS["panel2"], border=COLORS["cyan"])

    labels = ["ax", "ay", "az", "gx", "gy", "gz"]
    for idx, axis in enumerate(labels):
        lx = int(panel_left + Inches(0.38) + idx * Inches(1.75))
        add_textbox(slide, axis, lx, Inches(1.45), Inches(1.4), Inches(0.32), size_pt=24, bold=True, align=PP_ALIGN.CENTER)

    heat_left = int(panel_left + Inches(0.2))
    heat_top = int(panel_top + Inches(0.56))
    heat_w = int(panel_w - Inches(0.4))
    heat_h = int(panel_h - Inches(0.95))

    rng = random.Random(42)

    def big_values(r, c):
        t = r / 24
        if c < 3:
            v = 0.75 * math.sin(t * math.pi * 2.3 + c * 0.7) + 0.18 * math.cos(t * math.pi * 5 + c)
        else:
            v = 0.45 * math.sin(t * math.pi * 5.2 + c) + rng.uniform(-0.35, 0.35)
        return clamp(v, -1.0, 1.0)

    draw_heatmap(slide, heat_left, heat_top, heat_w, heat_h, rows=24, cols=6, value_fn=big_values)

    # Color legend
    legend_left = Inches(4.72)
    legend_top = Inches(4.33)
    legend_w = Inches(3.9)
    legend_h = Inches(0.2)
    for i in range(8):
        t = i / 7
        v = -1 + 2 * t
        block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(legend_left + legend_w * t),
            legend_top,
            int(legend_w / 8 + 1),
            legend_h,
        )
        fill_solid(block, heat_color(v))
        hide_line(block)
    add_textbox(slide, "low", Inches(4.2), Inches(4.24), Inches(0.5), Inches(0.32), size_pt=14, color=COLORS["muted"])
    add_textbox(slide, "high", Inches(8.65), Inches(4.24), Inches(0.6), Inches(0.32), size_pt=14, color=COLORS["muted"], align=PP_ALIGN.RIGHT)

    # Two comparison panels
    left_panel = add_card(slide, Inches(1.0), Inches(4.83), Inches(5.3), Inches(2.25), fill=COLORS["panel2"], border=(43, 113, 203))
    left_panel.adjustments[0] = 0.05
    right_panel = add_card(slide, Inches(7.03), Inches(4.83), Inches(5.3), Inches(2.25), fill=COLORS["panel2"], border=(43, 113, 203))
    right_panel.adjustments[0] = 0.05
    add_textbox(slide, "Wave", Inches(2.85), Inches(4.9), Inches(1.6), Inches(0.45), size_pt=30, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Shake", Inches(8.78), Inches(4.9), Inches(1.8), Inches(0.45), size_pt=30, bold=True, align=PP_ALIGN.CENTER)

    wave_rng = random.Random(7)
    shake_rng = random.Random(11)

    def wave_values(r, c):
        t = r / 16
        base = math.sin(t * math.pi * 2 + c * 0.9)
        return clamp(base * 0.8 + wave_rng.uniform(-0.1, 0.1), -1.0, 1.0)

    def shake_values(r, c):
        return clamp(shake_rng.uniform(-1.0, 1.0) * (0.6 + 0.4 * ((r + c) % 2)), -1.0, 1.0)

    draw_heatmap(slide, Inches(1.25), Inches(5.4), Inches(4.8), Inches(1.5), rows=16, cols=6, value_fn=wave_values)
    draw_heatmap(slide, Inches(7.28), Inches(5.4), Inches(4.8), Inches(1.5), rows=16, cols=6, value_fn=shake_values)

    add_note(
        slide,
        "Say this slowly: each gesture recording is 100 readings x 6 numbers = 600 numbers.\n"
        "Point to bottom heatmaps: Wave and Shake look like different number patterns.\n"
        "This is the core insight before students begin collecting data.",
    )


def build_slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, (20, 24, 92))
    add_glow(slide, Inches(2.0), Inches(1.2), Inches(9.4), Inches(4.8), COLORS["purple"], 0.65)
    add_glow(slide, Inches(4.15), Inches(2.1), Inches(4.2), Inches(2.7), COLORS["cyan"], 0.72)

    add_textbox(
        slide,
        "CHALLENGE TIME!",
        Inches(0.5),
        Inches(0.2),
        Inches(12.3),
        Inches(0.95),
        size_pt=58,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        "Can your AI score 7 out of 10?",
        Inches(2.7),
        Inches(1.05),
        Inches(7.9),
        Inches(0.55),
        size_pt=29,
        bold=True,
        color=COLORS["cyan"],
        align=PP_ALIGN.CENTER,
    )

    # Target rings
    centers = [(Inches(6.67), Inches(3.35), Inches(3.75), COLORS["cyan"], 0.85), (Inches(7.07), Inches(3.75), Inches(2.95), COLORS["purple"], 0.75), (Inches(7.47), Inches(4.15), Inches(2.15), COLORS["green"], 0.6)]
    for left, top, size, color, alpha in centers:
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        fill_solid(ring, color, alpha)
        set_line(ring, COLORS["white"], 1.4, transparency=0.65)

    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.53), Inches(2.2), Inches(2.25), Inches(2.25))
    fill_solid(center, (15, 46, 126))
    set_line(center, COLORS["cyan"], 2.0)
    add_textbox(slide, "7/10", Inches(5.64), Inches(2.86), Inches(2.05), Inches(0.6), size_pt=52, bold=True, align=PP_ALIGN.CENTER)

    goal = add_card(slide, Inches(4.45), Inches(4.95), Inches(4.45), Inches(0.8), fill=(33, 78, 158), border=COLORS["cyan"])
    goal.adjustments[0] = 0.08
    add_textbox(slide, "Goal: at least 7 correct attempts", Inches(4.5), Inches(5.06), Inches(4.3), Inches(0.48), size_pt=22, align=PP_ALIGN.CENTER, bold=True)

    steps = [("1", "See the target gesture"), ("2", "Do the move"), ("3", "Tap Score Attempt")]
    for idx, (num, text) in enumerate(steps):
        left = int(Inches(0.85) + idx * Inches(4.2))
        card = add_card(slide, left, Inches(6.02), Inches(3.86), Inches(1.18), fill=(17, 48, 126), border=(72, 146, 231))
        card.adjustments[0] = 0.07
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.22), Inches(6.24), Inches(0.58), Inches(0.58))
        fill_solid(badge, COLORS["yellow"])
        set_line(badge, COLORS["orange"], 1.2)
        add_textbox(slide, num, left + Inches(0.22), Inches(6.24), Inches(0.58), Inches(0.58), size_pt=20, bold=True, align=PP_ALIGN.CENTER, color=(45, 45, 45))
        add_textbox(slide, text, left + Inches(0.92), Inches(6.16), Inches(2.74), Inches(0.76), size_pt=18, bold=True, vertical_anchor=MSO_ANCHOR.TOP)

    add_note(
        slide,
        "Set the game target: 7 out of 10.\n"
        "Students should wait until predicted label matches before scoring each attempt.\n"
        "Focus: model quality comes from consistent movement and clear gesture differences.",
    )


def build_slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, (28, 35, 86))
    add_glow(slide, Inches(1.2), Inches(0.8), Inches(10.9), Inches(5.4), COLORS["blue"], 0.72)

    add_textbox(
        slide,
        "YOU JUST TRAINED A REAL NEURAL NETWORK!",
        Inches(0.35),
        Inches(0.32),
        Inches(12.6),
        Inches(0.9),
        size_pt=44,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Aptos Display",
    )

    lessons = [
        ("GOOD DATA\n=\nGOOD AI", COLORS["yellow"]),
        ("DIFFERENT\nPATTERNS HELP", COLORS["cyan"]),
        ("MORE PRACTICE,\nBETTER RESULTS", COLORS["green"]),
        ("AI IS MATH,\nNOT MAGIC", COLORS["pink"]),
    ]
    for idx, (label, accent) in enumerate(lessons):
        cx = Inches(0.75 + idx * 3.18)
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, Inches(1.62), Inches(2.65), Inches(2.65))
        fill_solid(bubble, (21, 63, 132))
        set_line(bubble, accent, 2.2)
        add_textbox(
            slide,
            label,
            cx + Inches(0.24),
            Inches(2.26),
            Inches(2.15),
            Inches(1.45),
            size_pt=18,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    band = add_card(slide, Inches(0.9), Inches(4.6), Inches(11.55), Inches(2.26), fill=(14, 42, 102), border=(68, 133, 219))
    band.adjustments[0] = 0.05
    add_textbox(slide, "Same idea, different scale", Inches(4.58), Inches(5.06), Inches(4.2), Inches(0.5), size_pt=29, bold=True, align=PP_ALIGN.CENTER, color=COLORS["cyan"])

    left_card = add_card(slide, Inches(1.24), Inches(5.2), Inches(2.95), Inches(1.24), fill=(17, 57, 126), border=(101, 167, 247))
    right_card = add_card(slide, Inches(9.18), Inches(5.2), Inches(2.95), Inches(1.24), fill=(17, 57, 126), border=(101, 167, 247))
    left_card.adjustments[0] = 0.08
    right_card.adjustments[0] = 0.08
    add_textbox(slide, "YOU TODAY\n30 samples\n6 sensors, 1 Arduino", Inches(1.42), Inches(5.26), Inches(2.6), Inches(1.05), size_pt=18, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "REAL AI\nMillions of samples\nHundreds of sensors", Inches(9.36), Inches(5.26), Inches(2.6), Inches(1.05), size_pt=18, bold=True, align=PP_ALIGN.CENTER)

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.54), Inches(5.45), Inches(3.2), Inches(0.7))
    fill_solid(arrow, COLORS["teal"])
    hide_line(arrow)

    add_note(
        slide,
        "Celebrate achievement first.\n"
        "Recap the 4 lessons from class.\n"
        "Close with transfer idea: same ML principles scale to robotics and self-driving systems.",
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    prs.save(OUTPUT)
    print(f"Saved deck: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
